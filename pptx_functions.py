"""
pptx_functions  v0.1.0
=======================
python-pptx abstraction layer for programmatic PowerPoint slide generation.

Provides a configuration-dictionary workflow for building slides without
writing python-pptx boilerplate.  Each slide object (text box, image,
connector, auto shape, table, banner, header) is driven by a plain Python
dict, making slide configurations easy to generate, serialise, and reuse.

Quick start
-----------
>>> from pptx_functions import create_slide_deck, get_default_config, OBJECT_TYPE_HANDLERS
>>> from datetime import datetime, timezone
>>> slide_config = {
...     "Details": {
...         "Author": "B. Rodriguez",
...         "Filename": "test.pptx",
...         "Slide Aspect Ratio": "16:9",
...         "Slide Width & Height": [13.33, 7.5],
...     },
...     "Slides": {
...         "Slide 01": {
...             "Slide Template": 6,
...             "Slide Name": "Test Slide",
...             "Slide Notes": "",
...             "Background Color": "#ffffff",
...             "Background Alpha": 0.0,
...             "Objects": {
...                 "Title": get_default_config("Title", {"Text": "Hello World"}),
...             },
...         }
...     },
... }

Functions
---------
Utilities
    get_image_details       Extract image metadata (size, DPI, format) using Pillow.
    color_to_rgb            Convert a CSS colour name or HEX string to (R, G, B).
    extract_ltwh            Return [left, top, width, height] as Inches from a config dict.

PPTX Object Functions
    create_slide_deck       Create a new Presentation with configured dimensions and metadata.
    set_slide_background    Apply background colour and transparency to a slide.
    add_autoshape           Add an auto shape to a slide.
    add_connector           Add a connector line to a slide.
    add_image               Add an image with smart aspect-ratio handling.
    add_notes               Write text to the Notes pane of a slide.
    add_shape_formatting    Apply line and fill formatting to any shape.
    add_table               Add a formatted table with optional per-cell styling.
    add_textbox             Add a text box (plain, multi-paragraph, multi-run, or bulleted).
    add_text_formatting     Apply font formatting to a text run or paragraph.

Wrapper Functions
    add_banners             Add top and bottom classification/marking banners to a slide.
    add_header              Add a slide header (title text box, connector, and seal images).

High-level API
    dispatch_objects             Render all enabled objects from an "Objects" config dict.
    build_presentation           Build a complete Presentation from a slide_config dict.
    extract_presentation_config  Extract a slide_config dict from an existing Presentation.

Helper Functions
    get_default_config      Return a deep-copied default config for a named object type.
    apply_metadata          Apply document metadata to a Presentation's core properties.
    show_functions          Print all public functions and data objects in this module.
    show_autoshapes         Print available auto shape key strings.
    show_object_alignment   Print available text alignment key strings.
    show_dash_styles        Print available line dash style key strings.
    show_slide_templates    Print slide layout template descriptions.

Data
    PPTX_LOOKUP             Consolidated enum lookup dict (align, valign, dash_styles,
                            shapes, connectors).
    OBJECT_TYPE_HANDLERS    Dispatch table mapping "Object Type" strings to functions.
    default_configurations  Default config dicts for every supported object type.
"""

from __future__ import annotations

__version__ = "0.2.0"
__author__  = "KChadwick"

__all__ = [
    # utilities
    "get_image_details", "color_to_rgb", "extract_ltwh",
    # pptx object functions
    "create_slide_deck", "set_slide_background",
    "add_autoshape", "add_connector", "add_image", "add_notes",
    "add_shape_formatting", "add_table", "add_textbox", "add_text_formatting",
    # wrapper functions
    "add_banners", "add_header",
    # high-level API
    "dispatch_objects", "build_presentation", "extract_presentation_config",
    # helper functions
    "get_default_config", "apply_metadata",
    "show_functions", "show_autoshapes", "show_object_alignment",
    "show_dash_styles", "show_slide_templates",
    # data
    "PPTX_LOOKUP", "OBJECT_TYPE_HANDLERS", "default_configurations",
]

import os
from copy import deepcopy
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple

from PIL import Image
import pptx
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.parts.coreprops import CorePropertiesPart as CoreProperties
from pptx.shapes.autoshape import Shape
from pptx.shapes.connector import Connector
from pptx.shapes.picture import Picture
from pptx.slide import Slide
from pptx.table import Table, _Cell
from pptx.util import Inches, Pt
import webcolors


################################################################################
# LOOKUP DATA
# Consolidated enum mappings used by all PPTX object functions.
################################################################################

PPTX_LOOKUP: Dict[str, Dict[str, Any]] = {

    "align": {
        "center":      PP_ALIGN.CENTER,
        "left":        PP_ALIGN.LEFT,
        "right":       PP_ALIGN.RIGHT,
        "justify":     PP_ALIGN.JUSTIFY,
        "justify_low": PP_ALIGN.JUSTIFY_LOW,
        "mixed":       PP_ALIGN.MIXED,
    },

    "valign": {
        "top":    MSO_VERTICAL_ANCHOR.TOP,
        "middle": MSO_VERTICAL_ANCHOR.MIDDLE,
        "bottom": MSO_VERTICAL_ANCHOR.BOTTOM,
    },

    "dash_styles": {
        "solid":         MSO_LINE_DASH_STYLE.SOLID,
        "-":             MSO_LINE_DASH_STYLE.SOLID,
        "dash":          MSO_LINE_DASH_STYLE.DASH,
        "--":            MSO_LINE_DASH_STYLE.DASH,
        "dash dot":      MSO_LINE_DASH_STYLE.DASH_DOT,
        "-.":            MSO_LINE_DASH_STYLE.DASH_DOT,
        "dash dot dot":  MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        "-..":           MSO_LINE_DASH_STYLE.DASH_DOT_DOT,
        "long dash":     MSO_LINE_DASH_STYLE.LONG_DASH,
        "long dash dot": MSO_LINE_DASH_STYLE.LONG_DASH_DOT,
        "round dot":     MSO_LINE_DASH_STYLE.ROUND_DOT,
        ".":             MSO_LINE_DASH_STYLE.ROUND_DOT,
        "square dot":    MSO_LINE_DASH_STYLE.SQUARE_DOT,
    },

    "shapes": {
        "trapezoid":         MSO_SHAPE.TRAPEZOID,
        "cube":              MSO_SHAPE.CUBE,
        "parallelogram":     MSO_SHAPE.PARALLELOGRAM,
        "sun":               MSO_SHAPE.SUN,
        "moon":              MSO_SHAPE.MOON,
        "lightning bolt":    MSO_SHAPE.LIGHTNING_BOLT,
        "hexagon":           MSO_SHAPE.HEXAGON,
        "heptagon":          MSO_SHAPE.HEPTAGON,
        "octagon":           MSO_SHAPE.OCTAGON,
        "cloud":             MSO_SHAPE.CLOUD,
        "gear 6":            MSO_SHAPE.GEAR_6,
        "gear 9":            MSO_SHAPE.GEAR_9,
        "explosion 1":       MSO_SHAPE.EXPLOSION1,
        "explosion 2":       MSO_SHAPE.EXPLOSION2,
        "smiley face":       MSO_SHAPE.SMILEY_FACE,
        "heart":             MSO_SHAPE.HEART,
        "rectangle":         MSO_SHAPE.RECTANGLE,
        "rounded rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
        "oval":              MSO_SHAPE.OVAL,
        "diamond":           MSO_SHAPE.DIAMOND,
        "right triangle":    MSO_SHAPE.RIGHT_TRIANGLE,
        "bent arrow":        MSO_SHAPE.BENT_ARROW,
        "right arrow":       MSO_SHAPE.RIGHT_ARROW,
        "left arrow":        MSO_SHAPE.LEFT_ARROW,
        "up arrow":          MSO_SHAPE.UP_ARROW,
        "down arrow":        MSO_SHAPE.DOWN_ARROW,
        "left right arrow":  MSO_SHAPE.LEFT_RIGHT_ARROW,
        "up down arrow":     MSO_SHAPE.UP_DOWN_ARROW,
        "star 4-point":      MSO_SHAPE.STAR_4_POINT,
        "star 5-point":      MSO_SHAPE.STAR_5_POINT,
        "star 6-point":      MSO_SHAPE.STAR_6_POINT,
        "star 7-point":      MSO_SHAPE.STAR_7_POINT,
        "star 8-point":      MSO_SHAPE.STAR_8_POINT,
        "star 16-point":     MSO_SHAPE.STAR_16_POINT,
        "star 32-point":     MSO_SHAPE.STAR_32_POINT,
        "chevron":           MSO_SHAPE.CHEVRON,
        "pentagon":          MSO_SHAPE.PENTAGON,
        "oval callout":      MSO_SHAPE.OVAL_CALLOUT,
        "rounded callout":   MSO_SHAPE.ROUNDED_RECTANGULAR_CALLOUT,
        "rectangular callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
        "cloud callout":     MSO_SHAPE.CLOUD_CALLOUT,
        "default":           MSO_SHAPE.RECTANGLE,
    },

    "connectors": {
        "straight": MSO_CONNECTOR.STRAIGHT,
        "elbow":    MSO_CONNECTOR.ELBOW,
        "curved":   MSO_CONNECTOR.CURVE,
    },
}

# Reverse lookups used by extract_presentation_config
_ALIGN_PREFERRED: Dict[Any, str] = {
    PP_ALIGN.CENTER:      "center",
    PP_ALIGN.LEFT:        "left",
    PP_ALIGN.RIGHT:       "right",
    PP_ALIGN.JUSTIFY:     "justify",
    PP_ALIGN.JUSTIFY_LOW: "justify_low",
}

_VALIGN_PREFERRED: Dict[Any, str] = {
    MSO_VERTICAL_ANCHOR.TOP:    "top",
    MSO_VERTICAL_ANCHOR.MIDDLE: "middle",
    MSO_VERTICAL_ANCHOR.BOTTOM: "bottom",
}

_DASH_PREFERRED: Dict[Any, str] = {
    MSO_LINE_DASH_STYLE.SOLID:         "-",
    MSO_LINE_DASH_STYLE.DASH:          "--",
    MSO_LINE_DASH_STYLE.DASH_DOT:      "-.",
    MSO_LINE_DASH_STYLE.DASH_DOT_DOT:  "-..",
    MSO_LINE_DASH_STYLE.LONG_DASH:     "long dash",
    MSO_LINE_DASH_STYLE.LONG_DASH_DOT: "long dash dot",
    MSO_LINE_DASH_STYLE.ROUND_DOT:     ".",
    MSO_LINE_DASH_STYLE.SQUARE_DOT:    "square dot",
}

_SHAPE_REVERSE: Dict[Any, str] = {
    v: k for k, v in PPTX_LOOKUP["shapes"].items() if k != "default"
}


_SLIDE_TEMPLATES: Dict[str, str] = {
    "0": "Title Slide — title and subtitle placeholders.",
    "1": "Title and Content — title at top, one content placeholder.",
    "2": "Section Header — single title placeholder.",
    "3": "Two Content — title at top, two side-by-side content placeholders.",
    "4": "Comparison — title, two content placeholders with headings.",
    "5": "Title Only — title placeholder, no content area.",
    "6": "Blank — no placeholders (most common for programmatic use).",
    "7": "Content with Caption — content, title, and caption placeholder.",
    "8": "Picture with Caption — picture, title, and caption placeholder.",
}


################################################################################
# UTILITIES
################################################################################

def get_image_details(img_path: str) -> Dict[str, Any]:
    """
    Extract metadata from an image file using Pillow.

    Parameters
    ----------
    img_path : str
        Path to the image file.  Supports all formats Pillow can read
        (JPEG, PNG, GIF, TIFF, WebP, BMP, ICO, etc.).

    Returns
    -------
    dict
        Keys: ``file_path``, ``file_format``, ``color_mode``, ``width_px``,
        ``height_px``, ``aspect_ratio``, ``size_pixels``, ``has_alpha``,
        ``is_animated``, ``n_frames``, ``metadata``, ``dpi``.

    Raises
    ------
    FileNotFoundError
        If *img_path* does not exist.
    """
    if not os.path.exists(img_path):
        raise FileNotFoundError(f"Image not found: {img_path!r}")

    with Image.open(img_path) as im:
        im.load()
        return {
            "file_path":    os.path.abspath(img_path),
            "file_format":  im.format,
            "color_mode":   im.mode,
            "width_px":     im.width,
            "height_px":    im.height,
            "aspect_ratio": im.width / im.height,
            "size_pixels":  im.size,
            "has_alpha":    im.mode in ("RGBA", "LA", "PA"),
            "is_animated":  getattr(im, "is_animated", False),
            "n_frames":     getattr(im, "n_frames", 1),
            "metadata":     dict(im.info),
            "dpi":          im.info.get("dpi", (72.0, 72.0)),
        }


def color_to_rgb(color: str) -> Tuple[int, int, int]:
    """
    Convert a CSS colour name or HEX code to an (R, G, B) integer tuple.

    Parameters
    ----------
    color : str
        A HEX string (e.g. ``"#cb181d"``) or CSS3 colour name (e.g. ``"red"``).

    Returns
    -------
    tuple of int
        ``(R, G, B)`` with each component in [0, 255].
        Returns ``(0, 0, 0)`` (black) and prints a warning for invalid input.

    Examples
    --------
    >>> color_to_rgb("#cb181d")
    (203, 24, 29)
    >>> color_to_rgb("steelblue")
    (70, 130, 180)
    """
    try:
        if color.startswith("#"):
            return webcolors.hex_to_rgb(color)
        return webcolors.name_to_rgb(color)
    except (ValueError, AttributeError):
        print(f"Warning: {color!r} is not a valid colour name or HEX code. Using black.")
        return (0, 0, 0)


def extract_ltwh(config: Dict[str, Any]) -> List:
    """
    Return ``[left, top, width, height]`` as ``Inches`` objects from *config*.

    Only keys present in *config* are included, so partial configs (e.g. a
    connector with only position keys) are handled correctly.

    Parameters
    ----------
    config : dict
        Must contain at minimum ``"left"`` and ``"top"``; ``"width"`` and
        ``"height"`` are included when present.

    Returns
    -------
    list of pptx.util.Inches
    """
    return [Inches(config[k]) for k in ("left", "top", "width", "height") if k in config]


################################################################################
# PPTX OBJECT FUNCTIONS
################################################################################

def _get_slide_size(slide: Slide) -> Tuple[float, float]:
    """Return (width_inches, height_inches) from the slide's parent Presentation."""
    prs = slide.part.package.presentation_part.presentation
    return prs.slide_width.inches, prs.slide_height.inches


def set_slide_background(slide: Slide, config: Dict[str, Any]) -> None:
    """
    Apply a solid background colour to a slide.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        ``"Background Color"`` (str, optional)
            HEX or CSS colour.  Default ``"#ffffff"``.
        ``"Background Alpha"`` (float, optional)
            Transparency 0.0 (opaque) – 1.0 (transparent).  Default ``0.0``.
    """
    r, g, b = color_to_rgb(config.get("Background Color", "#ffffff"))
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb          = RGBColor(r, g, b)
    fill.fore_color.transparency = config.get("Background Alpha", 0.0)


def create_slide_deck(
    config_details: Dict[str, Any],
    verbose: bool = False,
) -> Presentation:
    """
    Create a new PowerPoint presentation with configured dimensions and metadata.

    Sets the slide size, stamps the current UTC time as the creation date, and
    applies all recognised metadata fields.  Does not add any slides.

    Parameters
    ----------
    config_details : dict
        Presentation-level configuration.  Required keys:

        ``"Slide Width & Height"``
            ``[width_inches, height_inches]`` — use ``[13.33, 7.5]`` for 16:9
            or ``[10.0, 7.5]`` for 4:3.

        Optional keys (applied as core properties):
        ``"Author"``, ``"Title"``, ``"Subject"``, ``"Keywords"``,
        ``"Comments"``, ``"Category"``, ``"Description"``.

    verbose : bool, optional
        Print configuration details to stdout.  Default ``False``.

    Returns
    -------
    pptx.presentation.Presentation

    Examples
    --------
    >>> details = {"Slide Width & Height": [13.33, 7.5], "Author": "B. Rodriguez"}
    >>> prs = create_slide_deck(details)
    """
    prs = Presentation()
    config_details["Created"] = datetime.now(timezone.utc)

    slide_width, slide_height = config_details["Slide Width & Height"]
    prs.slide_width  = Inches(slide_width)
    prs.slide_height = Inches(slide_height)

    if verbose:
        print("Configuration Details:")
        for key, val in config_details.items():
            print(f"  {key}: {val}")
        print(f"\n  Slide Width:  {slide_width:.4f} inches")
        print(f"  Slide Height: {slide_height:.4f} inches")
        print(f"  Aspect Ratio: {slide_width / slide_height:.3f}:1\n")

    apply_metadata(prs.core_properties, config_details)
    return prs


def add_autoshape(slide: Slide, config: Dict[str, Any]) -> Shape:
    """
    Add an auto shape to a slide.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Configuration keys:

        ``"AutoShape Key"`` (str)
            Shape name — see :func:`show_autoshapes` for valid values.
            Defaults to ``"rectangle"``.
        ``"left"``, ``"top"``, ``"width"``, ``"height"`` (float)
            Position and size in inches.
        ``"Fill Color"`` (str, optional)
            HEX or CSS name for the shape fill.
        ``"Fill Alpha"`` (float, optional)
            Fill opacity 0.0–1.0.  Default ``1.0``.
        ``"Line Color"`` (str, optional)
            HEX or CSS name for the border.
        ``"Line Width"`` (float, optional)
            Border width in points.  Default ``0`` (no border).
        ``"Line Style"`` (str, optional)
            Dash style — see :func:`show_dash_styles`.  Default ``"-"`` (solid).

    Returns
    -------
    Shape
    """
    left, top, width, height = extract_ltwh(config)
    shape_type = PPTX_LOOKUP["shapes"].get(
        config.get("AutoShape Key", "default"), MSO_SHAPE.RECTANGLE
    )
    shape: Shape = slide.shapes.add_shape(
        autoshape_type_id=shape_type,
        left=left, top=top, width=width, height=height,
    )
    add_shape_formatting(shape, config)

    # Optional: render text directly into the shape's text frame
    if config.get("Text"):
        tf = shape.text_frame
        tf.word_wrap       = config.get("Word Wrap?", True)
        tf.vertical_anchor = PPTX_LOOKUP["valign"].get(
            config.get("V Align", "middle").lower(), MSO_VERTICAL_ANCHOR.MIDDLE
        )
        tf.clear()
        p           = tf.paragraphs[0]
        p.alignment = PPTX_LOOKUP["align"].get(
            config.get("Align", "center").lower(), PP_ALIGN.CENTER
        )
        run      = p.add_run()
        run.text = str(config["Text"])
        add_text_formatting(run, config)

    return shape


def add_connector(slide: Slide, config: Dict[str, Any]) -> Connector:
    """
    Add a connector line to a slide.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Configuration keys:

        ``"Line Color"`` (str)
            Line colour as HEX or CSS name.  Default ``"#000000"``.
        ``"Start X"``, ``"Start Y"``, ``"End X"``, ``"End Y"`` (float)
            Endpoint coordinates in inches.
        ``"Line Width"`` (float, optional)
            Line thickness in points.  Default ``1``.
        ``"Line Style"`` (str, optional)
            Dash style.  Default ``"-"`` (solid).
        ``"Type"`` (str, optional)
            ``"straight"``, ``"elbow"``, or ``"curved"``.  Default ``"straight"``.

    Returns
    -------
    Connector
    """
    # "Line Color/Width/Style" are the canonical keys; "Color/Width/Style" kept for compat
    rgb       = color_to_rgb(config.get("Line Color", config.get("Color", "#000000")))
    start_x   = Inches(config.get("Start X", 0))
    start_y   = Inches(config.get("Start Y", 0))
    end_x     = Inches(config.get("End X", 0))
    end_y     = Inches(config.get("End Y", 0))
    conn_type = PPTX_LOOKUP["connectors"].get(
        config.get("Type", "straight"), MSO_CONNECTOR.STRAIGHT
    )

    line: Connector = slide.shapes.add_connector(conn_type, start_x, start_y, end_x, end_y)
    line.line.color.rgb  = RGBColor(*rgb)
    line.line.width      = Pt(config.get("Line Width", config.get("Width", 1)))
    line.line.dash_style = PPTX_LOOKUP["dash_styles"].get(
        config.get("Line Style", config.get("Style", "-")), MSO_LINE_DASH_STYLE.SOLID
    )
    return line


def add_image(slide: Slide, config: Dict[str, Any]) -> Picture:
    """
    Add an image to a slide with smart aspect-ratio handling.

    Aspect ratio behaviour is controlled by ``"Preserve Aspect Ratio?"`` and
    ``"fit"``:

    ========================  ==================================================
    ``"fit"`` value           Behaviour
    ========================  ==================================================
    ``"width"``               Scale to the given width; compute height.
    ``"height"``              Scale to the given height; compute width.
    ``"native"`` (default)    Use the image's native pixel size at 96 DPI.
    ========================  ==================================================

    When ``"Preserve Aspect Ratio?"`` is ``False``, both ``"width"`` and
    ``"height"`` are applied exactly (stretching if needed).

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Required keys: ``"img_path"``, ``"left"``, ``"top"``.

        Optional keys: ``"width"``, ``"height"``, ``"fit"``,
        ``"Preserve Aspect Ratio?"``, ``"Flip Horizontal?"``,
        ``"Flip Vertical?"``, ``"Line Width"``, ``"Line Color"``,
        ``"Line Style"``.

    Returns
    -------
    Picture

    Raises
    ------
    FileNotFoundError
        If ``"img_path"`` does not point to an existing file.
    KeyError
        If any required key is missing from *config*.
    """
    for key in ("img_path", "left", "top"):
        if key not in config:
            raise KeyError(f"add_image: missing required config key '{key}'")

    left     = Inches(config["left"])
    top      = Inches(config["top"])
    details  = get_image_details(config["img_path"])
    aspect   = details["aspect_ratio"]
    native_w = details["width_px"] / 96.0
    native_h = details["height_px"] / 96.0

    preserve = config.get("Preserve Aspect Ratio?", True)
    fit_mode = config.get("fit", "native").lower()

    if preserve:
        if "width" in config and fit_mode in ("width", "native"):
            width  = Inches(config["width"])
            height = Inches(config["width"] / aspect)
        elif "height" in config and fit_mode == "height":
            height = Inches(config["height"])
            width  = Inches(config["height"] * aspect)
        else:
            width  = Inches(native_w)
            height = Inches(native_h)
    else:
        width  = Inches(config.get("width", native_w))
        height = Inches(config.get("height", native_h))

    picture: Picture = slide.shapes.add_picture(
        config["img_path"], left, top, width=width, height=height
    )

    # Support both Title Case ("Line Width") and legacy snake_case ("line_width")
    lw = config.get("Line Width", config.get("line_width", 0))
    if lw > 0:
        lc = config.get("Line Color", config.get("line_color", "#000000"))
        ls = config.get("Line Style", config.get("line_style", "-"))
        picture.line.color.rgb  = RGBColor(*color_to_rgb(lc))
        picture.line.width      = Pt(lw)
        picture.line.dash_style = PPTX_LOOKUP["dash_styles"].get(ls, MSO_LINE_DASH_STYLE.SOLID)

    # Flip / mirror
    if config.get("Flip Horizontal?", False) or config.get("Flip Vertical?", False):
        xfrm = picture.element.spPr.xfrm
        if config.get("Flip Horizontal?", False):
            xfrm.set("flipH", "1")
        if config.get("Flip Vertical?", False):
            xfrm.set("flipV", "1")

    return picture


def add_notes(slide: Slide, note_text: str) -> None:
    """
    Write text to the Notes pane of a slide.

    Parameters
    ----------
    slide : Slide
        Target slide.
    note_text : str
        Text to place in the notes pane.  Replaces any existing content.
    """
    slide.notes_slide.notes_text_frame.text = note_text


def add_shape_formatting(shape: Shape, config: Dict[str, Any]) -> Shape:
    """
    Apply line and fill formatting to a PowerPoint shape.

    Transparency is applied via an overlay rectangle workaround because
    python-pptx's ``fore_color.transparency`` is unreliable across
    PowerPoint versions.

    Parameters
    ----------
    shape : Shape
        The shape to format (AutoShape, TextBox, etc.).  Modified in place.
    config : dict
        Formatting keys:

        ``"Line Color"`` (str, optional)   Border colour.
        ``"Line Width"`` (float, optional) Border width in points; ``0`` = none.
        ``"Line Style"`` (str, optional)   Dash style.  Default ``"-"`` (solid).
        ``"Fill Color"`` (str, optional)   Fill colour.  ``None`` = transparent.
        ``"Fill Alpha"`` (float, optional) Opacity 0.0–1.0.  Default ``1.0``.

    Returns
    -------
    Shape
        The same shape object (modified in place).
    """
    # ── Line formatting ───────────────────────────────────────────────────────
    line_color = config.get("Line Color")
    line_width = config.get("Line Width", 0)
    if line_color and line_width > 0:
        shape.line.color.rgb  = RGBColor(*color_to_rgb(line_color))
        shape.line.width      = Pt(line_width)
        shape.line.dash_style = PPTX_LOOKUP["dash_styles"].get(
            config.get("Line Style", "-"), MSO_LINE_DASH_STYLE.SOLID
        )

    # ── Fill formatting ───────────────────────────────────────────────────────
    fill_color = config.get("Fill Color")
    if not fill_color:
        return shape

    r, g, b = color_to_rgb(fill_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)

    # Transparency via overlay rectangle (most reliable cross-version approach)
    alpha = config.get("Fill Alpha", 1.0)
    if 0.0 <= alpha < 0.999:
        try:
            slide = shape.part.slide
            overlay: Shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, shape.left, shape.top, shape.width, shape.height
            )
            overlay.fill.solid()
            overlay.fill.fore_color.rgb          = RGBColor(r, g, b)
            overlay.fill.fore_color.transparency = 1.0 - alpha
            overlay.line.fill.background()
        except Exception as exc:
            print(f"Warning: transparency overlay failed — {exc}")

    return shape


def add_table(slide: Slide, config: Dict[str, Any]) -> Table:
    """
    Add a formatted table to a slide.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Configuration keys:

        ``"left"``, ``"top"``, ``"width"``, ``"height"`` (float)
            Overall table position and size in inches.
        ``"Rows"`` (int)
            Number of data rows (header row is added automatically).
        ``"Columns"`` (int)
            Number of columns.
        ``"Column Widths"`` (list of float)
            Width of each column in inches.
        ``"Row Height"`` (float)
            Height of each row in inches.  Default ``0.4``.
        ``"Column Headers"`` (list of str)
            Labels for the top row.
        ``"Row Data"`` (list of list)
            2-D list of cell values (str).
        ``"Align"`` (str, optional)
            Horizontal text alignment.  Default ``"left"``.
        ``"V-Align"`` (str, optional)
            Vertical cell alignment.  Default ``"middle"``.
        ``"Font Size"`` (int, optional)
            Font size in points.  Default ``10``.
        ``"Font Color"`` (str, optional)
            HEX or CSS colour for cell text.
        ``"Bold?"`` (bool, optional)
            ``True`` bolds all cells; header row is bold by default.
        ``"Cell Styles"`` (dict, optional)
            Per-cell formatting overrides.  Keys are ``(row_index, col_index)``
            tuples (0-based, header row = 0).  Values are config dicts with any
            combination of ``"Font Color"``, ``"Font Size"``, ``"Bold?"``,
            ``"V-Align"``,
            ``"Italic?"``, ``"Underline?"``, ``"Fill Color"``, ``"Align"``.

            Example::

                "Cell Styles": {
                    (0, 2): {"Font Color": "#ffffff", "Fill Color": "#1a1a2e"},
                    (2, 3): {"Font Color": "#cb181d", "Bold?": True},
                }

    Returns
    -------
    Table
    """
    left, top, width, height = extract_ltwh(config)
    rows_count = config.get("Rows", 0) + 1
    cols_count = config.get("Columns", 0)

    shape        = slide.shapes.add_table(rows_count, cols_count, left, top, width, height)
    table: Table = shape.table

    for idx, col_w in enumerate(config.get("Column Widths", [])):
        if idx < len(table.columns):
            table.columns[idx].width = Inches(col_w)

    for row in table.rows:
        row.height = Inches(config.get("Row Height", 0.4))

    align_key   = config.get("Align", "left").lower()
    valign_key  = config.get("V-Align", "middle").lower()
    full_data   = [config.get("Column Headers", [])] + config.get("Row Data", [])
    cell_styles = config.get("Cell Styles", {})

    for row_idx, row_values in enumerate(full_data):
        for col_idx, cell_value in enumerate(row_values):
            if col_idx >= cols_count:
                break
            cell: _Cell = table.cell(row_idx, col_idx)

            # Merge base config with any per-cell overrides
            override    = cell_styles.get((row_idx, col_idx), {})
            merged      = {**config, "Bold?": config.get("Bold?", row_idx == 0), **override}

            cell.vertical_anchor = PPTX_LOOKUP["valign"].get(
                merged.get("V-Align", valign_key).lower(), MSO_VERTICAL_ANCHOR.MIDDLE
            )
            cell_align  = PPTX_LOOKUP["align"].get(
                merged.get("Align", align_key).lower(), PP_ALIGN.LEFT
            )

            # Cell background fill
            if "Fill Color" in override:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*color_to_rgb(override["Fill Color"]))

            tf = cell.text_frame
            tf.clear()
            p           = tf.paragraphs[0]
            p.alignment = cell_align
            run         = p.add_run()
            run.text    = str(cell_value)
            add_text_formatting(run, merged)

    return table


def add_textbox(slide: Slide, config: Dict[str, Any]) -> Shape:
    """
    Add a text box to a slide.

    Supports four text modes selected automatically from the config:

    1. **Bulleted list** — when ``"Bullets"`` is present.
       Value: ``[[text, level], ...]`` where level 0 = top level.
       ``"Font Size"`` may be a scalar or a ``[[level, size], ...]`` list.

    2. **Multi-paragraph** — when ``"Text"`` is a ``list``.
       Each element is one paragraph: a ``str`` for plain text, or a
       ``dict`` of ``{run_text: run_config}`` for mixed formatting.

    3. **Multi-run text** — when ``"Text"`` is a ``dict``.
       Value: ``{run_text: {per-run overrides}, ...}``.
       Useful for mixed formatting within one paragraph (e.g. bold + colour).

    4. **Plain text** — when ``"Text"`` is a ``str`` (default).

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Configuration keys:

        ``"left"``, ``"top"``, ``"width"``, ``"height"`` (float)
            Position and size in inches.
        ``"Text"`` (str or dict)
            Text content.  Required unless ``"Bullets"`` is provided.
        ``"Bullets"`` (list of [str, int], optional)
            Overrides ``"Text"``; drives bulleted-list mode.
        ``"Align"`` (str, optional)
            Horizontal alignment.  Default ``"center"``.
        ``"V Align"`` (str, optional)
            Vertical alignment.  Default ``"middle"``.
        ``"Font Name"`` (str, optional)    Default ``"Calibri"``.
        ``"Font Size"`` (int or list)      Default ``18``.
        ``"Font Color"`` (str, optional)   Default ``"#535353"``.
        ``"Bold?"``, ``"Italic?"``, ``"Underline?"`` (bool, optional)
        ``"Word Wrap?"`` (bool, optional)  Default ``True``.
        ``"Fill Color"``, ``"Fill Alpha"`` — passed to :func:`add_shape_formatting`.
        ``"Line Color"``, ``"Line Width"``, ``"Line Style"`` — border options.

    Returns
    -------
    Shape
        The text box shape object.
    """
    left, top, width, height = extract_ltwh(config)
    tx_box: Shape = slide.shapes.add_textbox(left, top, width, height)
    add_shape_formatting(tx_box, config)

    tf = tx_box.text_frame
    tf.word_wrap      = config.get("Word Wrap?", True)
    tf.auto_size      = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = PPTX_LOOKUP["valign"].get(
        config.get("V Align", "middle").lower(), MSO_VERTICAL_ANCHOR.MIDDLE
    )
    tf.clear()

    align_key = config.get("Align", "center").lower()

    # tf.clear() leaves one empty paragraph behind; reuse it as the first
    # paragraph rather than calling add_paragraph(), which would create a
    # second paragraph and produce a blank leading line in every text box.
    def _first_para():
        return tf.paragraphs[0]

    def _next_para():
        return tf.add_paragraph()

    # ── Bulleted list ─────────────────────────────────────────────────────────
    if "Bullets" in config:
        font_size_raw = config.get("Font Size", 12)
        if isinstance(font_size_raw, (list, tuple)):
            font_size_map     = dict(font_size_raw)
            default_font_size = 12
        elif isinstance(font_size_raw, dict):
            font_size_map     = font_size_raw
            default_font_size = 12
        else:
            font_size_map     = {}
            default_font_size = font_size_raw

        for i, (bullet_text, level) in enumerate(config["Bullets"]):
            p           = _first_para() if i == 0 else _next_para()
            p.level     = level
            p.alignment = PPTX_LOOKUP["align"].get(align_key, PP_ALIGN.LEFT)
            run         = p.add_run()
            run.text    = bullet_text
            font_size   = font_size_map.get(level, default_font_size)
            add_text_formatting(run, {**config, "Font Size": font_size})

    # ── Multi-paragraph ───────────────────────────────────────────────────────
    elif isinstance(config.get("Text"), list):
        for i, para_content in enumerate(config["Text"]):
            p           = _first_para() if i == 0 else _next_para()
            p.alignment = PPTX_LOOKUP["align"].get(align_key, PP_ALIGN.CENTER)
            if isinstance(para_content, dict):
                for text, run_config in para_content.items():
                    run      = p.add_run()
                    run.text = text
                    add_text_formatting(run, {**config, **run_config})
            else:
                run      = p.add_run()
                run.text = str(para_content)
                add_text_formatting(run, config)

    # ── Multi-run text (single paragraph) ────────────────────────────────────
    elif isinstance(config.get("Text"), dict):
        p           = _first_para()
        p.alignment = PPTX_LOOKUP["align"].get(align_key, PP_ALIGN.CENTER)
        for text, run_config in config["Text"].items():
            run      = p.add_run()
            run.text = text
            add_text_formatting(run, {**config, **run_config})

    # ── Plain text ────────────────────────────────────────────────────────────
    else:
        p           = _first_para()
        p.alignment = PPTX_LOOKUP["align"].get(align_key, PP_ALIGN.CENTER)
        run         = p.add_run()
        run.text    = str(config.get("Text", ""))
        add_text_formatting(run, config)

    return tx_box


def add_text_formatting(run, config: Dict[str, Any]) -> None:
    """
    Apply font formatting to a text run or paragraph.

    Parameters
    ----------
    run : pptx run or paragraph
        The object whose ``.font`` attribute will be set.
    config : dict
        Formatting keys:

        ``"Font Name"`` (str)    Font family.  Default ``"Calibri"``.
        ``"Font Size"`` (float)  Size in points.  Default ``12``.
        ``"Font Color"`` (str)   HEX or CSS colour.  Default ``"#000000"``.
        ``"Bold?"`` (bool)       Default ``False``.
        ``"Italic?"`` (bool)     Default ``False``.
        ``"Underline?"`` (bool)  Default ``False``.
    """
    run.font.name      = config.get("Font Name", "Calibri")
    run.font.size      = Pt(config.get("Font Size", 12))
    run.font.bold      = config.get("Bold?", False)
    run.font.italic    = config.get("Italic?", False)
    run.font.underline = config.get("Underline?", False)

    font_color = config.get("Font Color", "#000000")
    if font_color:
        run.font.color.rgb = RGBColor(*color_to_rgb(font_color))


################################################################################
# WRAPPER FUNCTIONS
################################################################################

def add_banners(slide: Slide, config: Dict[str, Any]) -> None:
    """
    Add top and bottom classification/marking banners to a slide.

    Places two identical text boxes — one flush with the top edge, one flush
    with the bottom — spanning the full slide width.  Banner height and all
    text formatting are drawn from ``default_configurations["Banner"]`` and
    can be overridden via *config*.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Keys:

        ``"Text"`` (str)
            Banner text (e.g. ``"UNCLASSIFIED"``).

        All other Text config keys (``"Font Name"``, ``"Font Color"``, etc.)
        are merged on top of the Banner defaults.  Slide width is read
        automatically from the slide object.
    """
    slide_width, slide_height = _get_slide_size(slide)
    banner       = deepcopy(default_configurations["Banner"])
    banner.update({k: v for k, v in config.items() if k != "Add?"})
    banner["width"] = slide_width

    for top in (0.0, slide_height - banner["height"]):
        cfg       = deepcopy(banner)
        cfg["top"] = top
        add_textbox(slide, cfg)


def add_header(slide: Slide, config: Dict[str, Any]) -> None:
    """
    Add a slide header consisting of a title text box, a separator connector,
    and optional left and right seal images.

    Parameters
    ----------
    slide : Slide
        Target slide.
    config : dict
        Keys:

        ``"Text"`` (str)
            Title text for the header.
        ``"Slide Aspect Ratio"`` (str, optional)
            ``"16:9"`` or ``"4:3"``.  Default ``"4:3"``.
        ``"Header Connector"`` (dict, optional)
            Connector config dict.  ``"End X"`` is computed automatically
            from the slide width and ``"Start X"``.
        ``"Left Seal"`` (dict, optional)
            Image config dict for the left seal/logo.
        ``"Right Seal"`` (dict, optional)
            Image config dict for the right seal/logo.
        ``"Add Right Image?"`` (bool, optional)
            Set ``False`` to suppress the right seal.  Default ``True``.
    """
    slide_width, _ = _get_slide_size(slide)

    title_cfg         = deepcopy(default_configurations["Title"])
    title_cfg["Text"] = config.get("Text", "")
    title_cfg["width"] = slide_width - 1.5
    add_textbox(slide, title_cfg)

    if "Header Connector" in config:
        conn_cfg        = deepcopy(config["Header Connector"])
        conn_cfg["End X"] = slide_width - conn_cfg.get("Start X", 0)
        add_connector(slide, conn_cfg)

    if "Left Seal" in config:
        add_image(slide, config["Left Seal"])

    if config.get("Add Right Image?", True) and "Right Seal" in config:
        add_image(slide, config["Right Seal"])


################################################################################
# HIGH-LEVEL API
################################################################################

def dispatch_objects(slide: Slide, objects_config: Dict[str, Any]) -> None:
    """
    Render all enabled objects from an ``"Objects"`` config dict onto a slide.

    Iterates over every entry in *objects_config*, skips any with
    ``"Add?": False``, looks up the handler for ``"Object Type"``, and calls it.

    Parameters
    ----------
    slide : Slide
        Target slide.
    objects_config : dict
        Mapping of arbitrary object names to object config dicts, e.g.::

            {
                "Title": {"Add?": True, "Object Type": "Text", ...},
                "Logo":  {"Add?": True, "Object Type": "Image", ...},
                "Draft": {"Add?": False, "Object Type": "AutoShape", ...},
            }
    """
    for _name, elem_config in objects_config.items():
        if not elem_config.get("Add?", True):
            continue
        obj_type = elem_config.get("Object Type", "")
        handler  = OBJECT_TYPE_HANDLERS.get(obj_type)
        if handler:
            handler(slide, elem_config)


def build_presentation(
    slide_config: Dict[str, Any],
    verbose: bool = False,
) -> Presentation:
    """
    Build a complete PowerPoint presentation from a ``slide_config`` dict.

    This is the top-level entry point.  It calls :func:`create_slide_deck`,
    iterates over every slide definition, applies the background, renders all
    enabled objects via :func:`dispatch_objects`, and writes slide notes.

    Parameters
    ----------
    slide_config : dict
        Full presentation config with the following structure::

            {
                "Details": {
                    "Author":               "...",
                    "Title":                "...",
                    "Filename":             "output.pptx",
                    "Slide Aspect Ratio":   "16:9",
                    "Slide Width & Height": [13.33, 7.5],
                },
                "Slides": {
                    "Slide 01": {
                        "Slide Template":   6,
                        "Slide Name":       "Cover",
                        "Slide Notes":      "Presenter notes.",
                        "Background Color": "#ffffff",
                        "Background Alpha": 0.0,
                        "Objects": { ... },
                    },
                },
            }

    verbose : bool, optional
        Passed through to :func:`create_slide_deck`.  Default ``False``.

    Returns
    -------
    pptx.presentation.Presentation
        The fully built presentation.  Call ``.save(filename)`` to write it.

    Examples
    --------
    >>> prs = build_presentation(slide_config, verbose=True)
    >>> prs.save(slide_config["Details"]["Filename"])
    """
    config_details = slide_config["Details"]
    prs = create_slide_deck(config_details, verbose=verbose)

    for slide_name, slide_cfg in slide_config.get("Slides", {}).items():
        layout = prs.slide_layouts[slide_cfg.get("Slide Template", 6)]
        slide  = prs.slides.add_slide(layout)

        set_slide_background(slide, slide_cfg)

        note_text = slide_cfg.get("Slide Notes", "")
        if note_text:
            add_notes(slide, note_text)

        dispatch_objects(slide, slide_cfg.get("Objects", {}))

    return prs


################################################################################
# EXTRACTION HELPERS  (private)
################################################################################

_EMU_PER_INCH: int = 914400


def _safe_hex(color_obj) -> Optional[str]:
    """Return '#rrggbb' from a pptx color object, or None if the type is not RGB."""
    try:
        return f"#{str(color_obj.rgb).lower()}"
    except Exception:
        return None


def _emu_to_in(emu) -> float:
    """Convert EMU to inches, rounded to 4 decimal places."""
    if emu is None:
        return 0.0
    return round(emu / _EMU_PER_INCH, 4)


def _get_layout_index(slide: Slide) -> int:
    """Return the slide layout index within its slide master, defaulting to 6 (Blank)."""
    try:
        sl = slide.slide_layout
        for i, layout in enumerate(sl.slide_master.slide_layouts):
            if layout._element is sl._element:
                return i
    except Exception:
        pass
    return 6


def _extract_background(slide: Slide) -> Tuple[str, float]:
    """Return (hex_color, alpha) from the slide's explicit background fill."""
    try:
        fill = slide.background.fill
        if fill.type is not None:
            c = _safe_hex(fill.fore_color)
            if c:
                alpha = 0.0
                try:
                    alpha = float(fill.fore_color.transparency or 0.0)
                except Exception:
                    pass
                return c, alpha
    except Exception:
        pass
    return "#ffffff", 0.0


def _extract_notes_text(slide: Slide) -> str:
    """Return stripped notes text from a slide."""
    try:
        return slide.notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def _extract_font(font) -> Dict[str, Any]:
    """Pull explicitly-set font properties into a partial config dict."""
    out: Dict[str, Any] = {}
    if font.name:
        out["Font Name"] = font.name
    if font.size:
        try:
            out["Font Size"] = round(font.size.pt, 1)
        except Exception:
            pass
    if font.bold is not None:
        out["Bold?"] = bool(font.bold)
    if font.italic is not None:
        out["Italic?"] = bool(font.italic)
    if font.underline is not None:
        out["Underline?"] = bool(font.underline)
    c = _safe_hex(font.color)
    if c:
        out["Font Color"] = c
    return out


def _extract_line_props(shape) -> Dict[str, Any]:
    """Return Line Color/Width/Style from a shape's line properties."""
    out: Dict[str, Any] = {"Line Color": None, "Line Width": 0, "Line Style": "-"}
    try:
        line = shape.line
        c = _safe_hex(line.color)
        if c:
            out["Line Color"] = c
        if line.width:
            out["Line Width"] = round(line.width.pt, 2)
        if line.dash_style is not None:
            out["Line Style"] = _DASH_PREFERRED.get(line.dash_style, "-")
    except Exception:
        pass
    return out


def _extract_fill_props(shape) -> Dict[str, Any]:
    """Return Fill Color/Alpha from a shape's fill properties."""
    out: Dict[str, Any] = {"Fill Color": None, "Fill Alpha": 1.0}
    try:
        fill = shape.fill
        if fill.type is not None:
            c = _safe_hex(fill.fore_color)
            if c:
                out["Fill Color"] = c
    except Exception:
        pass
    return out


def _extract_text_config(shape) -> Dict[str, Any]:
    """Extract a Text object config from any shape with a text frame."""
    tf = shape.text_frame

    font_props = {
        "Font Name": "Calibri", "Font Size": 12.0, "Font Color": "#535353",
        "Bold?": False, "Italic?": False, "Underline?": False,
    }
    align   = "left"
    v_align = "middle"

    try:
        v_anchor = tf.vertical_anchor
        if v_anchor is not None:
            v_align = _VALIGN_PREFERRED.get(v_anchor, "middle")
    except Exception:
        pass

    for para in tf.paragraphs:
        if para.alignment:
            align = _ALIGN_PREFERRED.get(para.alignment, "left")
        font = para.runs[0].font if para.runs else para.font
        font_props.update(_extract_font(font))
        if para.text.strip():
            break

    text = "\n".join(p.text for p in tf.paragraphs).strip()

    config: Dict[str, Any] = {
        "Add?":        True,
        "Object Type": "Text",
        "Text":        text,
        "left":   _emu_to_in(shape.left),
        "top":    _emu_to_in(shape.top),
        "width":  _emu_to_in(shape.width),
        "height": _emu_to_in(shape.height),
        "Align":      align,
        "V Align":    v_align,
        "Word Wrap?": bool(tf.word_wrap) if tf.word_wrap is not None else True,
    }
    config.update(font_props)
    config.update(_extract_fill_props(shape))
    config.update(_extract_line_props(shape))
    return config


def _extract_image_config(shape) -> Dict[str, Any]:
    """Extract an Image object config; img_path is null (content placeholder)."""
    return {
        "Add?":                   True,
        "Object Type":            "Image",
        "img_path":               None,
        "Preserve Aspect Ratio?": True,
        "fit":                    "width",
        "left":   _emu_to_in(shape.left),
        "top":    _emu_to_in(shape.top),
        "width":  _emu_to_in(shape.width),
        "height": _emu_to_in(shape.height),
    }


def _extract_connector_config(shape) -> Dict[str, Any]:
    """Extract a Connector object config, accounting for flip orientation."""
    left  = _emu_to_in(shape.left)
    top   = _emu_to_in(shape.top)
    right = _emu_to_in(shape.left + shape.width)
    bot   = _emu_to_in(shape.top  + shape.height)

    start_x, end_x = left, right
    start_y, end_y = top,  bot
    try:
        xfrm = shape.element.spPr.xfrm
        if xfrm.flipH:
            start_x, end_x = end_x, start_x
        if xfrm.flipV:
            start_y, end_y = end_y, start_y
    except Exception:
        pass

    lp = _extract_line_props(shape)
    return {
        "Add?":        True,
        "Object Type": "Connector",
        "Type":        "straight",
        "Start X": start_x, "Start Y": start_y,
        "End X":   end_x,   "End Y":   end_y,
        "Line Color": lp.get("Line Color") or "#000000",
        "Line Width": lp.get("Line Width", 1),
        "Line Style": lp.get("Line Style", "-"),
    }


def _extract_autoshape_config(shape) -> Dict[str, Any]:
    """Extract an AutoShape object config, including text if present."""
    try:
        shape_key = _SHAPE_REVERSE.get(shape.auto_shape_type, "rectangle")
    except Exception:
        shape_key = "rectangle"

    config: Dict[str, Any] = {
        "Add?":          True,
        "Object Type":   "AutoShape",
        "AutoShape Key": shape_key,
        "left":   _emu_to_in(shape.left),
        "top":    _emu_to_in(shape.top),
        "width":  _emu_to_in(shape.width),
        "height": _emu_to_in(shape.height),
    }
    config.update(_extract_fill_props(shape))
    config.update(_extract_line_props(shape))

    if shape.has_text_frame and shape.text_frame.text.strip():
        tf = shape.text_frame
        font_props = {
            "Font Name": "Calibri", "Font Size": 12.0, "Font Color": "#535353",
            "Bold?": False, "Italic?": False,
        }
        align   = "center"
        v_align = "middle"
        try:
            v_anchor = tf.vertical_anchor
            if v_anchor is not None:
                v_align = _VALIGN_PREFERRED.get(v_anchor, "middle")
        except Exception:
            pass
        for para in tf.paragraphs:
            if para.alignment:
                align = _ALIGN_PREFERRED.get(para.alignment, "center")
            font = para.runs[0].font if para.runs else para.font
            font_props.update(_extract_font(font))
            if para.text.strip():
                break
        config["Text"]      = "\n".join(p.text for p in tf.paragraphs).strip()
        config["Align"]     = align
        config["V Align"]   = v_align
        config["Word Wrap?"] = bool(tf.word_wrap) if tf.word_wrap is not None else True
        config.update(font_props)

    return config


def _extract_table_config(shape) -> Dict[str, Any]:
    """Extract a Table object config including headers and row data."""
    table   = shape.table
    n_rows  = len(table.rows)
    n_cols  = len(table.columns)

    col_widths = [_emu_to_in(col.width) for col in table.columns]
    row_height = _emu_to_in(table.rows[0].height) if table.rows else 0.4

    all_data = [
        [cell.text_frame.text.strip() for cell in row.cells]
        for row in table.rows
    ]
    headers  = all_data[0]      if all_data           else []
    row_data = all_data[1:]     if len(all_data) > 1  else []

    font_props = {
        "Font Name": "Calibri", "Font Size": 10.0,
        "Font Color": "#333333", "Bold?": False,
    }
    align = "center"
    try:
        para = table.cell(0, 0).text_frame.paragraphs[0]
        if para.alignment:
            align = _ALIGN_PREFERRED.get(para.alignment, "center")
        font = para.runs[0].font if para.runs else para.font
        font_props.update(_extract_font(font))
    except Exception:
        pass

    config: Dict[str, Any] = {
        "Add?":           True,
        "Object Type":    "Table",
        "left":   _emu_to_in(shape.left),
        "top":    _emu_to_in(shape.top),
        "width":  _emu_to_in(shape.width),
        "height": _emu_to_in(shape.height),
        "Rows":           max(0, n_rows - 1),
        "Columns":        n_cols,
        "Column Widths":  col_widths,
        "Row Height":     row_height,
        "Column Headers": headers,
        "Row Data":       row_data,
        "Align":          align,
        "V-Align":        "middle",
    }
    config.update(font_props)
    return config


def _extract_objects(slide: Slide) -> Dict[str, Any]:
    """Extract all non-placeholder shapes from a slide into an Objects config dict."""
    from pptx.shapes.connector import Connector as _Conn
    from pptx.shapes.picture   import Picture   as _Pic

    objects: Dict[str, Any] = {}
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            continue
        name = shape.name
        if isinstance(shape, _Pic):
            objects[name] = _extract_image_config(shape)
        elif isinstance(shape, _Conn):
            objects[name] = _extract_connector_config(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            objects[name] = _extract_table_config(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            objects[name] = _extract_autoshape_config(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            objects[name] = _extract_text_config(shape)
    return objects


def extract_presentation_config(prs) -> Dict[str, Any]:
    """
    Extract a slide_config dict from an existing Presentation.

    Captures all layout and formatting properties.  Image ``img_path`` values
    are set to ``None`` — they are content, not structure — making the result
    ready to use as a template: inject new paths (and update ``"Text"``
    strings) then pass to :func:`build_presentation`.

    Parameters
    ----------
    prs : Presentation or str
        An opened python-pptx Presentation object, or a file path string to a
        ``.pptx`` file which will be opened automatically.

    Returns
    -------
    dict
        A ``slide_config`` dict compatible with :func:`build_presentation`.

    Examples
    --------
    >>> from pptx_functions import extract_presentation_config, build_presentation
    >>> template = extract_presentation_config("existing.pptx")
    >>> template["Slides"]["Slide 01"]["Objects"]["Hero Image"]["img_path"] = "new.png"
    >>> prs = build_presentation(template)
    >>> prs.save("output.pptx")
    """
    if isinstance(prs, str):
        prs = Presentation(prs)
    cp      = prs.core_properties
    slide_w = round(prs.slide_width.inches,  4)
    slide_h = round(prs.slide_height.inches, 4)
    ratio   = slide_w / slide_h
    if   abs(ratio - 16 / 9) < 0.02:  aspect = "16:9"
    elif abs(ratio - 4  / 3) < 0.02:  aspect = "4:3"
    else:                              aspect = f"{ratio:.3f}:1"

    config: Dict[str, Any] = {
        "Details": {
            "Author":               cp.author   or "",
            "Title":                cp.title    or "",
            "Subject":              cp.subject  or "",
            "Comments":             cp.comments or "",
            "Keywords":             cp.keywords or "",
            "Category":             cp.category or "",
            "Filename":             "extracted_template.pptx",
            "Slide Aspect Ratio":   aspect,
            "Slide Width & Height": [slide_w, slide_h],
        },
        "Slides": {},
    }

    for idx, slide in enumerate(prs.slides, start=1):
        key                  = f"Slide {idx:02d}"
        bg_color, bg_alpha   = _extract_background(slide)
        config["Slides"][key] = {
            "Slide Template":   _get_layout_index(slide),
            "Slide Name":       key,
            "Slide Notes":      _extract_notes_text(slide),
            "Background Color": bg_color,
            "Background Alpha": bg_alpha,
            "Objects":          _extract_objects(slide),
        }

    return config


################################################################################
# HELPER FUNCTIONS
################################################################################

def get_default_config(
    object_type: str,
    overrides: Optional[Dict[str, Any]] = None,
) -> Dict[str, Any]:
    """
    Return a deep-copied default configuration for a named object type.

    Parameters
    ----------
    object_type : str
        One of: ``"AutoShape"``, ``"Banner"``, ``"Connector"``, ``"Image"``,
        ``"Table"``, ``"Text"``, ``"Title"``.
    overrides : dict, optional
        Keys to update on top of the defaults.  The defaults are not mutated.

    Returns
    -------
    dict

    Raises
    ------
    ValueError
        If *object_type* is not a recognised key.

    Examples
    --------
    >>> cfg = get_default_config("Text", {"Text": "Hello", "Font Size": 24})
    >>> cfg["Object Type"]
    'Text'
    """
    if object_type not in default_configurations:
        raise ValueError(
            f"Unknown object type: {object_type!r}. "
            f"Supported types: {list(default_configurations.keys())}"
        )
    config = deepcopy(default_configurations[object_type])
    if overrides:
        config.update(overrides)
    return config


def apply_metadata(
    core_props: CoreProperties,
    metadata_dict: Dict[str, Any],
) -> None:
    """
    Apply document metadata to a Presentation's core properties.

    Recognised keys: ``"Author"``, ``"Category"``, ``"Comments"``,
    ``"Created"``, ``"Description"``, ``"Keywords"``,
    ``"Last Modified By"``, ``"Modified"``, ``"Subject"``, ``"Title"``.

    Unrecognised keys (e.g. ``"Filename"``, ``"Slide Aspect Ratio"``) are
    silently ignored — they are workflow config, not OOXML metadata.

    Parameters
    ----------
    core_props : CoreProperties
        ``prs.core_properties`` from a Presentation object.
    metadata_dict : dict
        Metadata key→value mapping.
    """
    _FIELDS = {
        "Author":          "author",
        "Category":        "category",
        "Comments":        "comments",
        "Created":         "created",
        "Description":     "description",
        "Keywords":        "keywords",
        "Last Modified By":"last_modified_by",
        "Modified":        "modified",
        "Subject":         "subject",
        "Title":           "title",
    }
    for key, value in metadata_dict.items():
        attr = _FIELDS.get(key)
        if attr and hasattr(core_props, attr):
            setattr(core_props, attr, value)


def show_functions() -> None:
    """Print all public functions and data objects in this module."""
    sections = {
        "Utilities": [
            "get_image_details(img_path)",
            "color_to_rgb(color)",
            "extract_ltwh(config)",
        ],
        "PPTX Object Functions": [
            "create_slide_deck(config_details, verbose=False)",
            "set_slide_background(slide, config)",
            "add_autoshape(slide, config)",
            "add_connector(slide, config)",
            "add_image(slide, config)",
            "add_notes(slide, note_text)",
            "add_shape_formatting(shape, config)",
            "add_table(slide, config)",
            "add_textbox(slide, config)",
            "add_text_formatting(run, config)",
        ],
        "Wrapper Functions": [
            "add_banners(slide, config)",
            "add_header(slide, config)",
        ],
        "High-level API": [
            "dispatch_objects(slide, objects_config)",
            "build_presentation(slide_config, verbose=False)",
            "extract_presentation_config(prs)",
        ],
        "Helper Functions": [
            "get_default_config(object_type, overrides=None)",
            "apply_metadata(core_props, metadata_dict)",
            "show_functions()",
            "show_autoshapes()",
            "show_object_alignment()",
            "show_dash_styles()",
            "show_slide_templates()",
        ],
        "Data": [
            "PPTX_LOOKUP          — enum lookup dict (align, valign, dash_styles, shapes, connectors)",
            "OBJECT_TYPE_HANDLERS — dispatch table for slide object rendering",
            "default_configurations — default config dicts for every object type",
        ],
    }
    print(f"\npptx_functions  v{__version__}")
    print("=" * 60)
    for section, items in sections.items():
        print(f"\n{section}:")
        for item in items:
            print(f"  - {item}")
    print('\nType help(<function_name>) for full parameter documentation.')


def show_autoshapes() -> None:
    """Print available auto shape key strings for use in ``"AutoShape Key"``."""
    print("Available AutoShape Keys:")
    for key in PPTX_LOOKUP["shapes"]:
        if key != "default":
            print(f"  - {key}")


def show_object_alignment() -> None:
    """Print available alignment key strings for use in ``"Align"``."""
    print("Available Alignment Keys:")
    for key in PPTX_LOOKUP["align"]:
        print(f"  - {key}")


def show_dash_styles() -> None:
    """Print available dash style key strings for use in ``"Line Style"``."""
    print("Available Dash Style Keys:")
    seen: set = set()
    for key, val in PPTX_LOOKUP["dash_styles"].items():
        if val not in seen:
            print(f"  - {key!r}")
            seen.add(val)


def show_slide_templates() -> None:
    """Print slide layout template descriptions."""
    print("Slide Templates:")
    for key, desc in _SLIDE_TEMPLATES.items():
        print(f"  {key}: {desc}")


################################################################################
# DEFAULT CONFIGURATIONS
################################################################################

default_configurations: Dict[str, Dict[str, Any]] = {

    "AutoShape": {
        "Add?":          True,
        "Object Type":   "AutoShape",
        "AutoShape Key": "rectangle",
        "left": 3.0,  "top": 3.0, "width": 4.0, "height": 2.5,
        "Fill Color":  "#6A0DAD", "Fill Alpha": 1.0,
        "Line Color":  "#FFA500", "Line Width": 2.5, "Line Style": "dash",
    },

    "Banner": {
        "Add?":        True,
        "Object Type": "Banner",
        "Text":        "",
        "left": 0.0,  "top": 0.0, "width": 10.0, "height": 0.4,
        "Align":       "center",
        "V Align":     "middle",
        "Font Name":   "Calibri",
        "Font Size":   14,
        "Font Color":  "#535353",
        "Bold?": False, "Italic?": False, "Underline?": False, "Word Wrap?": False,
        "Fill Color":  None,      "Fill Alpha": 1.0,
        "Line Color":  None,      "Line Width": 0, "Line Style": "-",
    },

    "Connector": {
        "Add?":        False,
        "Object Type": "Connector",
        "Type":        "straight",
        "Start X": 1.5, "Start Y": 1.0, "End X": 8.5, "End Y": 1.0,
        "Line Color": "#535353", "Line Width": 2, "Line Style": "-",
    },

    "Image": {
        "Add?":        True,
        "Object Type": "Image",
        "img_path":    "",
        "Preserve Aspect Ratio?": True,
        "fit":         "width",
        "left": 3.25, "top": 1.5, "width": 5.0, "height": 5.0,
        "Line Width": 0, "Line Color": "#000000", "Line Style": "-",
    },

    "Table": {
        "Add?":           True,
        "Object Type":    "Table",
        "left": 1.0, "top": 1.5, "width": 8.0, "height": 4.0,
        "Columns":        3,
        "Rows":           1,
        "Column Widths":  [2.0, 3.0, 3.0],
        "Row Height":     0.4,
        "Column Headers": ["Column 1", "Column 2", "Column 3"],
        "Row Data":       [],
        "Font Size":      10,
        "Font Color":     "#333333",
        "Align":          "center",
        "V-Align":        "middle",
        "Bold?":          False,
    },

    "Text": {
        "Add?":        True,
        "Object Type": "Text",
        "Text":        "",
        "left": 1.5, "top": 1.0, "width": 10.0, "height": 4.0,
        "Align":       "center",
        "V Align":     "middle",
        "Font Name":   "Calibri",
        "Font Size":   18,
        "Font Color":  "#535353",
        "Bold?": False, "Italic?": False, "Underline?": False, "Word Wrap?": True,
        "Fill Color":  None, "Fill Alpha": 1.0,
        "Line Color":  None, "Line Width": 0, "Line Style": "-",
    },

    "Title": {
        "Add?":        True,
        "Object Type": "Text",
        "Text":        "Slide Title",
        "left": 1.5, "top": 0.3, "width": 8.5, "height": 0.5,
        "Align":       "left",
        "V Align":     "middle",
        "Font Name":   "Calibri",
        "Font Size":   28,
        "Font Color":  "#535353",
        "Bold?": False, "Italic?": False, "Underline?": False, "Word Wrap?": True,
        "Fill Color":  None, "Fill Alpha": 1.0,
        "Line Color":  None, "Line Width": 0, "Line Style": "-",
    },
}


################################################################################
# DISPATCH TABLE
################################################################################

OBJECT_TYPE_HANDLERS: Dict[str, Any] = {
    "AutoShape": add_autoshape,
    "Banner":    add_banners,
    "Connector": add_connector,
    "Header":    add_header,
    "Image":     add_image,
    "Table":     add_table,
    "Text":      add_textbox,
}
